"""
按 5 段结构生成答辩 PPT（深色科技风背景）。
- 5 段：项目成果介绍 / 工作参与形式 / 技术难点及解决措施 / AI 工具使用说明 / 实习整体复盘
- 每段：扉页 + 2 页内容 = 3 页
- 1 封面 + 5×3 + 1 结尾 = 17 页
- 背景：scripts/_ppt_bg/*.png
- 统一大字号，无脚本/方法论细节，左下角不再放重叠字样
"""
import os
import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---------- 路径 ----------
ROOT = r"c:\Users\82028\Desktop\github网站\githubpage网站\awaqwq233.github.io"
BG_DIR = os.path.join(ROOT, "scripts", "_ppt_bg")
OUT = os.path.join(ROOT, "scripts", "固定收益部_数据库性能优化实践_李颜铭_2026.pptx")

# ---------- 配色 ----------
WHITE = RGBColor(0xF5, 0xF8, 0xFF)
SOFT = RGBColor(0xCB, 0xD5, 0xE5)
DIM = RGBColor(0x8A, 0x9A, 0xB5)
DARK = RGBColor(0x10, 0x1A, 0x2E)

# 5 段 accent
ACCENT = {
    1: RGBColor(0x00, 0xC8, 0xE0),   # 青蓝
    2: RGBColor(0x90, 0xD8, 0x5A),   # 绿
    3: RGBColor(0xF5, 0xA8, 0x42),   # 橙
    4: RGBColor(0xC4, 0x82, 0xE8),   # 紫
    5: RGBColor(0xF5, 0x7A, 0x90),   # 粉
}

# ---------- 幻灯片 ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- 工具 ----------
def add_bg(slide, png_name):
    """把背景 PNG 铺满整页"""
    slide.shapes.add_picture(os.path.join(BG_DIR, png_name),
                             0, 0, prs.slide_width, prs.slide_height)


def set_text(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             font="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    tf.clear()
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = 1.20
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, text, size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb.text_frame, text, size, bold, color, align, font, anchor)
    return tb


def add_multiline(slide, x, y, w, h, lines, size=18, color=WHITE,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.35, font="Microsoft YaHei",
                  default_bold=False):
    """lines: list of dict {text, size?, bold?, color?, space_after?}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get("align", align)
        p.line_spacing = item.get("line_spacing", line_spacing)
        if "space_after" in item:
            p.space_after = Pt(item["space_after"])
        text = item.get("text", "")
        # 允许 runs：[(text, size, bold, color), ...]
        runs = item.get("runs")
        if runs is None:
            runs = [(text, item.get("size", size), item.get("bold", default_bold), item.get("color", color))]
        for j, (txt, s, b, c) in enumerate(runs):
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(s)
            r.font.bold = b
            r.font.color.rgb = c
            rPr = r._r.get_or_add_rPr()
            ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", font)
    return tb


def add_line(slide, x1, y1, x2, y2, color, width=1.5):
    cn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(width)
    return cn


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.fill.fore_color.brightness = 0  # keep
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


# ---------- 页脚：右上角小角标，不与正文重叠 ----------
def add_corner_badge(slide, idx, total, accent):
    """右上角：PART x / 5 + 进度小条"""
    add_text(slide, 10.5, 0.25, 2.5, 0.4,
             f"PART {idx} / 5  ·  {total} P",
             size=11, color=accent, align=PP_ALIGN.RIGHT, bold=True)
    # 进度条
    bar_w = 2.0
    bar_x = 13.333 - 0.5 - bar_w
    add_rect(slide, bar_x, 0.72, bar_w, 0.04, fill=DIM)
    add_rect(slide, bar_x, 0.72, bar_w * idx / 5, 0.04, fill=accent)


def add_brand(slide):
    """右下角品牌水印（不占左下角避免与正文重叠）"""
    add_text(slide, 8.5, 7.15, 4.5, 0.3,
             "固定收益部 · 数据库性能优化实践 · 2026",
             size=10, color=DIM, align=PP_ALIGN.RIGHT)


# ============================================================
# 1) 封面
# ============================================================
def build_cover():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, "01_cover.png")
    # 顶部小条
    add_rect(s, 0, 0, 13.333, 0.05, fill=ACCENT[1])
    # 主标
    add_text(s, 0.8, 1.4, 11.5, 0.6,
             "INTERNSHIP DEFENSE  ·  2026",
             size=14, color=ACCENT[1], bold=True, font="Consolas")
    add_text(s, 0.8, 2.1, 12.0, 1.3,
             "固定收益部 · 数据库性能优化实践",
             size=48, color=WHITE, bold=True)
    add_text(s, 0.8, 3.4, 12.0, 0.7,
             "Fixed Income · Database Performance Optimization",
             size=20, color=SOFT, font="Consolas")
    # 分隔线
    add_line(s, 0.8, 4.4, 6.5, 4.4, ACCENT[1], 2)
    # 介绍
    add_multiline(s, 0.8, 4.7, 12.0, 1.6, [
        {"text": "汇报人：李颜铭", "size": 22, "color": WHITE, "bold": True, "space_after": 6},
        {"text": "指导部门：固定收益部 · 风控与数据团队", "size": 18, "color": SOFT, "space_after": 4},
        {"text": "学校：武汉大学 · 计算机学院 · 软件工程 2028 届", "size": 18, "color": SOFT, "space_after": 4},
        {"text": "实习时间：2026-07-08 ~ 2026-07-28（3 周）", "size": 18, "color": SOFT},
    ])
    # 右下角版本号
    add_text(s, 10.5, 7.15, 2.5, 0.3,
             "V 2.0  ·  PPT for Defense",
             size=10, color=DIM, align=PP_ALIGN.RIGHT, font="Consolas")
    add_brand(s)


# ============================================================
# 2) Section 扉页（5 个）
# ============================================================
def build_divider(idx, title_en, total_pages):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, f"sec{idx}_divider.png")
    accent = ACCENT[idx]
    # 角标
    add_corner_badge(s, idx, total_pages, accent)
    # 段编号
    add_text(s, 0.8, 2.4, 12.0, 0.6,
             f"PART  0{idx}",
             size=22, color=accent, bold=True, font="Consolas")
    # 主标
    add_text(s, 0.8, 3.1, 12.0, 1.2,
             title_en,
             size=56, color=WHITE, bold=True)
    # 副标（中文小字）
    sub_map = {
        1: "Project Overview  ·  课题 · 交付物 · 目标与实现",
        2: "Work Participation  ·  团队结构 · 个人角色 · 核心产出",
        3: "Technical Challenges  ·  排查思路 · 解决路径 · 落地结果",
        4: "AI Tool Usage  ·  真实场景 · 使用边界 · 体感复盘",
        5: "Internship Review  ·  收获 · 成长 · 待提升方向",
    }
    add_text(s, 0.8, 4.4, 12.0, 0.5,
             sub_map[idx],
             size=18, color=SOFT, font="Consolas")
    # 大段分隔
    add_line(s, 0.8, 5.1, 4.0, 5.1, accent, 3)
    add_brand(s)


# ============================================================
# 3) Section 内容页
# 每段 2 页。每页统一布局：左 60% 大段文字，右 40% 一张"卡片图"
# ============================================================
def build_content(idx, page_no, title, blocks, total_pages):
    """
    blocks: list of dict
       { "type": "kpi" | "para" | "list" | "flow",  ... }
    """
    s = prs.slides.add_slide(BLANK)
    add_bg(s, f"sec{idx}_content.png")
    accent = ACCENT[idx]

    # 顶部小标签
    add_text(s, 0.8, 0.35, 8.0, 0.4,
             f"PART 0{idx}  ·  PAGE {page_no} / 2",
             size=12, color=accent, bold=True, font="Consolas")
    add_corner_badge(s, idx, total_pages, accent)

    # 标题
    add_text(s, 0.8, 0.8, 11.5, 0.9,
             title,
             size=34, color=WHITE, bold=True)
    # 标题下短分隔
    add_line(s, 0.8, 1.78, 2.4, 1.78, accent, 3)

    # 内容区
    content_x = 0.8
    content_y = 2.05
    content_w = 11.5
    content_h = 4.9

    for b in blocks:
        kind = b.get("type")
        x = content_x + b.get("x", 0)
        y = content_y + b.get("y", 0)
        w = b.get("w", content_w)
        h = b.get("h", 1.0)
        if kind == "para":
            add_multiline(s, x, y, w, h, b["lines"], size=b.get("size", 20),
                          color=b.get("color", WHITE),
                          line_spacing=b.get("line_spacing", 1.45),
                          default_bold=b.get("bold", False))
        elif kind == "list":
            # 自定义项目符号
            lines = []
            for it in b["items"]:
                if isinstance(it, str):
                    lines.append({"runs": [
                        ("▍  ", 18, True, accent),
                        (it, b.get("size", 19), False, b.get("color", SOFT)),
                    ], "line_spacing": 1.55, "space_after": 8})
                else:
                    runs = [("▍  ", 18, True, accent)]
                    for seg in it:
                        runs.append(tuple(seg))
                    lines.append({"runs": runs, "line_spacing": 1.55, "space_after": 8})
            add_multiline(s, x, y, w, h, lines, size=b.get("size", 19),
                          color=b.get("color", SOFT), line_spacing=1.55)
        elif kind == "kpi":
            # 数据卡：底色 + accent 左边条
            card_x, card_y, card_w, card_h = x, y, w, h
            add_rect(s, card_x, card_y, card_w, card_h,
                     fill=RGBColor(0x14, 0x22, 0x40), line=accent, line_w=1.2)
            add_rect(s, card_x, card_y, 0.08, card_h, fill=accent)
            add_multiline(s, card_x + 0.25, card_y + 0.12, card_w - 0.4, card_h - 0.2,
                          [
                              {"text": b["big"], "size": b.get("big_size", 32), "color": accent, "bold": True, "space_after": 4},
                              {"text": b["label"], "size": b.get("label_size", 15), "color": WHITE, "bold": True, "space_after": 2},
                              {"text": b.get("sub", ""), "size": b.get("sub_size", 12), "color": SOFT},
                          ], size=20, line_spacing=1.25)
        elif kind == "flow":
            # 横向流程节点
            n = len(b["steps"])
            step_w = w / n
            for i, st in enumerate(b["steps"]):
                cx = x + step_w * i
                add_rect(s, cx + 0.15, y, step_w - 0.3, h, fill=RGBColor(0x14, 0x22, 0x40), line=accent, line_w=0.8)
                add_text(s, cx + 0.15, y + 0.15, step_w - 0.3, 0.5,
                         st["no"], size=14, color=accent, bold=True, font="Consolas",
                         align=PP_ALIGN.CENTER)
                add_text(s, cx + 0.15, y + 0.55, step_w - 0.3, h - 0.6,
                         st["label"], size=st.get("size", 16), color=WHITE, bold=True,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                if i < n - 1:
                    add_text(s, cx + step_w - 0.18, y, 0.36, h,
                             "▶", size=18, color=accent, bold=True,
                             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif kind == "box":
            # 通用卡片
            add_rect(s, x, y, w, h, fill=RGBColor(0x14, 0x22, 0x40), line=accent, line_w=1.0)
            add_multiline(s, x + 0.25, y + 0.15, w - 0.4, h - 0.2, b["lines"],
                          size=b.get("size", 18), color=b.get("color", WHITE),
                          line_spacing=1.5)
    add_brand(s)


# ============================================================
# 4) 结尾页
# ============================================================
def build_ending():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, "99_ending.png")
    add_text(s, 0.8, 2.5, 12.0, 0.6,
             "THANKS FOR LISTENING",
             size=20, color=ACCENT[1], bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 3.2, 12.0, 1.4,
             "感谢聆听",
             size=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 4.7, 12.0, 0.6,
             "Q & A  ·  欢迎各位老师与同学指正",
             size=22, color=SOFT, align=PP_ALIGN.CENTER)
    add_line(s, 5.5, 5.7, 7.8, 5.7, ACCENT[1], 2)
    add_text(s, 0.8, 6.0, 12.0, 0.5,
             "固定收益部 · 数据库性能优化实践  ·  李颜铭  ·  2026",
             size=14, color=DIM, align=PP_ALIGN.CENTER, font="Consolas")


# ============================================================
# 内容数据
# ============================================================
TOTAL = 16  # 1 cover + 5 sec * (1 divider + 2 content) ≈ 16 pages


# ---------- Section 1：项目成果介绍 ----------
def sec1_pages():
    # P1: 课题 + 目标
    build_content(1, 1, "课题背景与目标", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 1.1,
         "lines": [
             {"text": "固定收益部结算链路长期受数据库性能瓶颈制约，", "size": 22, "color": WHITE, "space_after": 8},
             {"text": "本次实习以可量化的性能指标为锚点，完成核心过程的优化与上线验证。", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.4, "w": 5.85, "h": 3.4,
         "lines": [
             {"text": "课题 · 三件事", "size": 22, "color": ACCENT[1], "bold": True, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("压测基线  ", 20, True, WHITE),
                       ("覆盖核心结算链路，刻画优化前后的真实差距", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("数据对账  ", 20, True, WHITE),
                       ("账务结果零漂移，是优化上线的硬门槛", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("监控能力  ", 20, True, WHITE),
                       ("把一次性脚本沉淀为可持续观察的指标", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 1.4, "w": 5.85, "h": 3.4,
         "lines": [
             {"text": "目标 · 三个数字", "size": 22, "color": ACCENT[1], "bold": True, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("核心 P95 耗时显著下降", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("账务数据零漂移可证明", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 12},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("形成可复用的工程方法", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)

    # P2: 交付物 + 个人负责模块
    build_content(1, 2, "交付物与个人模块", [
        {"type": "kpi", "x": 0, "y": 0, "w": 3.85, "h": 1.7,
         "big": "LEVERAGE", "big_size": 32, "label": "并发压测",
         "sub": "建立阶梯压测基线，刻画核心链路并发行为", "label_size": 17, "sub_size": 13},
        {"type": "kpi", "x": 4.05, "y": 0, "w": 3.85, "h": 1.7,
         "big": "D3", "big_size": 32, "label": "数据核对",
         "sub": "账务数据零漂移核对，确保优化可上线", "label_size": 17, "sub_size": 13},
        {"type": "kpi", "x": 8.1, "y": 0, "w": 3.85, "h": 1.7,
         "big": "Dashboard", "big_size": 32, "label": "指标看板",
         "sub": "把零散结果汇总为可持续观察的指标视图", "label_size": 17, "sub_size": 13},
        {"type": "box", "x": 0, "y": 2.0, "w": 12.0, "h": 2.7,
         "lines": [
             {"text": "我负责的模块", "size": 22, "color": ACCENT[1], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("LEVERAGE 压测基线  ", 20, True, WHITE),
                       ("— 端到端脚本 + 跨档位执行 + 报表产出", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("D3 数据核对  ", 20, True, WHITE),
                       ("— 多源对账 + 差异定位 + 一致性结论", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[1]),
                       ("指标看板  ", 20, True, WHITE),
                       ("— 把脚本产出转化为团队可读的指标视图", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)


# ---------- Section 2：工作参与形式 ----------
def sec2_pages():
    # P1: 团队结构
    build_content(2, 1, "团队结构与个人角色", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "所在小组由导师、后端开发、DBA 与数据四个角色构成，", "size": 22, "color": WHITE, "space_after": 6},
             {"text": "我以实习生的身份独立承担压测与对账两条主线。", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 12.0, "h": 1.5,
         "lines": [
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("团队构成  ", 20, True, WHITE),
                       ("导师  ·  后端  ·  DBA  ·  数据", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("我的角色  ", 20, True, WHITE),
                       ("独立承担压测基线 + 数据对账两条主线，需求与上线由团队共同 review", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 0, "y": 2.7, "w": 5.85, "h": 2.1,
         "lines": [
             {"text": "我做的事", "size": 22, "color": ACCENT[2], "bold": True, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("脚本编写、脚本运行、报告整理", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("对账方案设计、异常样本复核", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("把一次性的结果汇总为团队可读的指标", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 2.7, "w": 5.85, "h": 2.1,
         "lines": [
             {"text": "上游给我的", "size": 22, "color": ACCENT[2], "bold": True, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("需求背景  ·  数据范围  ·  上线节奏", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("测试环境  ·  必要的运维协助", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("关键样本与结论的最终 review", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 4},
         ]},
    ], TOTAL)

    # P2: 协作流程
    build_content(2, 2, "协作流程与核心产出", [
        {"type": "flow", "x": 0, "y": 0, "w": 12.0, "h": 1.3,
         "steps": [
             {"no": "STEP 01", "label": "需求对齐"},
             {"no": "STEP 02", "label": "方案设计"},
             {"no": "STEP 03", "label": "实施验证"},
             {"no": "STEP 04", "label": "团队 review"},
         ]},
        {"type": "box", "x": 0, "y": 1.7, "w": 12.0, "h": 3.0,
         "lines": [
             {"text": "我的核心产出", "size": 22, "color": ACCENT[2], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("可复用的压测脚本与对账脚本  ", 20, True, WHITE),
                       ("不依赖我本人在场也能继续运行", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("可读的性能报告  ", 20, True, WHITE),
                       ("把数据 + 结论 + 边界写在同一份文档里", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[2]),
                       ("可承接的工程方法  ", 20, True, WHITE),
                       ("在团队 review 后能直接被下一任接手", 19, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)


# ---------- Section 3：技术难点及解决措施 ----------
def sec3_pages():
    # P1: 难点 1
    build_content(3, 1, "难点一 · 并发压测场景的设计", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "从单点压测到全链路压测，需要回答：", "size": 22, "color": WHITE, "space_after": 6},
             {"text": "“线上真实负载下，瓶颈到底在哪里？”", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "排查思路", "size": 22, "color": ACCENT[3], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("从单接口压测开始，先找局部最弱点", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("逐步叠加链路节点，定位耦合处", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("用真实数据分布替代均匀造数", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("以 P95 而非平均值刻画真实表现", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "落地结果", "size": 22, "color": ACCENT[3], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("建立可复用的阶梯压测基线", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("在多档位下都能跑出稳定结论", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("压测与对账结果进入看板，被团队复用", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)

    # P2: 难点 2 + 通用化沉淀
    build_content(3, 2, "难点二 · 数据一致性的核对", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "性能优化的硬门槛不是更快，", "size": 22, "color": WHITE, "space_after": 6},
             {"text": "而是账务结果不能出现任何漂移。", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "核对思路", "size": 22, "color": ACCENT[3], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("同一份真实输入跑优化前 / 后两套逻辑", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("逐字段比较输出，差异立即定位", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("对异常样本人工复核，不掩盖", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("结果先于性能，结论写在文档里", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "通用化沉淀", "size": 22, "color": ACCENT[3], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("把一次性的脚本沉淀为可复用机制", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("降低下次同类优化的边际成本", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[3]),
                       ("形成团队可读、可 review 的工程产出", 19, True, WHITE)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)


# ---------- Section 4：AI 工具使用说明 ----------
def sec4_pages():
    # P1: 应用场景
    build_content(4, 1, "我在哪些环节用 AI", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "AI 是效率放大器，不是决策替代品。", "size": 22, "color": WHITE, "space_after": 6},
             {"text": "我把它放在重复劳动与上下文切换最多的地方。", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 5.85, "h": 1.7,
         "lines": [
             {"text": "脚本与代码", "size": 22, "color": ACCENT[4], "bold": True, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("样板代码、参数解析、异常分支", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("快速出第一版，再人工改造", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 1.0, "w": 5.85, "h": 1.7,
         "lines": [
             {"text": "文档与汇报", "size": 22, "color": ACCENT[4], "bold": True, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("周报草稿、复盘结构、表达润色", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("关键事实仍以我自己的笔记为准", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 4},
         ]},
        {"type": "box", "x": 0, "y": 2.9, "w": 12.0, "h": 2.0,
         "lines": [
             {"text": "思路与对比", "size": 22, "color": ACCENT[4], "bold": True, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("在排查复杂问题时，让 AI 列多种切入角度", 18, False, SOFT)],
              "line_spacing": 1.5, "space_after": 6},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("我再筛选并亲自动手验证", 19, True, WHITE)],
              "line_spacing": 1.5, "space_after": 4},
         ]},
    ], TOTAL)

    # P2: 使用边界
    build_content(4, 2, "使用边界与体会", [
        {"type": "box", "x": 0, "y": 0, "w": 5.85, "h": 4.7,
         "lines": [
             {"text": "我用得顺的场景", "size": 22, "color": ACCENT[4], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("脚本脚手架、参数模板", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("周报与复盘的结构化表达", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("把同一份事实改写成不同读者", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("复杂问题的多角度切入", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 0, "w": 5.85, "h": 4.7,
         "lines": [
             {"text": "我不放心交给 AI 的", "size": 22, "color": ACCENT[4], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("业务口径与数据范围的最终判断", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("性能数字、生产影响、上线策略", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("账务一致性的最终结论", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[4]),
                       ("对客户、对外部的承诺性表述", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)


# ---------- Section 5：实习整体复盘 ----------
def sec5_pages():
    # P1: 收获与成长
    build_content(5, 1, "三个月的实习收获", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "把零散的脚本经验，沉淀成一套能跨场景复用的工程方法。", "size": 22, "color": WHITE},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 3.85, "h": 3.8,
         "lines": [
             {"text": "方法", "size": 22, "color": ACCENT[5], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("从基线到结论的完整链路", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("以 P95 而非均值做判断", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("结果先于性能", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 4.05, "y": 1.0, "w": 3.85, "h": 3.8,
         "lines": [
             {"text": "能力", "size": 22, "color": ACCENT[5], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("独立设计可复用的脚本", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("在团队中把结论讲清楚", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("把过程沉淀为可交接的产出", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 8.1, "y": 1.0, "w": 3.85, "h": 3.8,
         "lines": [
             {"text": "态度", "size": 22, "color": ACCENT[5], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("把每一个数字写到能 review", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("承认不知道的，比假装知道重要", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("复盘比结论本身更有价值", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)

    # P2: 待提升方向
    build_content(5, 2, "待提升方向与相关思考", [
        {"type": "para", "x": 0, "y": 0, "w": 12.0, "h": 0.8,
         "lines": [
             {"text": "知道边界在哪，比已经会什么更重要。", "size": 22, "color": WHITE, "space_after": 6},
             {"text": "以下是我给自己列的下一阶段方向。", "size": 22, "color": SOFT},
         ]},
        {"type": "box", "x": 0, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "技术深度", "size": 22, "color": ACCENT[5], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("更系统地吃透执行计划与锁机制", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("在更真实的数据分布下做稳定性压测", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("把 AWR / SQL Trace 等工具串成日常工作流", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
        {"type": "box", "x": 6.15, "y": 1.0, "w": 5.85, "h": 3.8,
         "lines": [
             {"text": "工程与协作", "size": 22, "color": ACCENT[5], "bold": True, "space_after": 10},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("把脚本沉淀为团队长期维护的工具", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("补齐生产灰度监控与回归基线", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 8},
             {"runs": [("▍  ", 18, True, ACCENT[5]),
                       ("在跨角色协作中保持节奏感与边界感", 18, False, SOFT)],
              "line_spacing": 1.55, "space_after": 4},
         ]},
    ], TOTAL)


# ============================================================
# 装配
# ============================================================
build_cover()
sec1_pages()
build_divider(2, "工作参与形式", TOTAL)
sec2_pages()
build_divider(3, "技术难点及解决措施", TOTAL)
sec3_pages()
build_divider(4, "AI 工具使用说明", TOTAL)
sec4_pages()
build_divider(5, "实习整体复盘", TOTAL)
sec5_pages()
build_ending()

prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
